from pptx import Presentation
from pptx.util import Inches
import datetime
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import pandas as pd
import numpy as np
import datetime as dt
import dataframe_image as dfi
import os

class RevenueMakerClass:
    def __init__(self, include_support_cost=False, filter_by_year=False, switch_product_group=False, column1=None, column2=None):
        print("Initializing RevenueMakerClass...")
        self.include_support_cost = include_support_cost
        self.filter_by_year = filter_by_year
        self.switch_product_group = switch_product_group
        self.column1 = column1
        self.column2 = column2
        self.image_dir = os.path.join(os.getcwd(), 'image_bin')

        if not os.path.exists(self.image_dir):
            os.makedirs(self.image_dir)

        self.support_cost = pd.read_csv('.csv')
        self.product_categorization = pd.read_csv('.csv')
        #self.revenue_data = pd.read_csv('revenue_data.csv')
        self.revenue_data = pd.read_csv('.csv')


        ### provisory revenue processing
        self.revenue_data = self.revenue_data[self.revenue_data['Product Group'] != 'ONLY_ABCD']
        self.revenue_data = self.revenue_data.drop(columns=['Support Cost'], errors='ignore')

        ### Year filter & Group / Type filter
        if self.filter_by_year:
            self.revenue_data = self.revenue_data[(self.revenue_data['Year'] == 2023) | (self.revenue_data['Year'] == 2024)]
        if self.switch_product_group:
            self.revenue_data.rename(columns={'Product Group': 'Product Type', 'Product Type': 'Product Group'}, inplace=True)
            self.support_cost.rename(columns={'Product Group': 'Product Type', 'Product Type': 'Product Group'}, inplace=True)
            self.product_categorization.rename(columns={'Product Group': 'Product Type', 'Product Type': 'Product Group'}, inplace=True)

    def compute_ONLY_ABCD(self):
        product_types = ["Product_A", "Product_B", "Product_C", "Product_D"]
        filtered_df = self.revenue_data[self.revenue_data['Product Type'].isin(product_types)]
        filtered_df['Revenue'] = pd.to_numeric(filtered_df['Revenue'], errors='coerce')
        filtered_df['Revenue'] = filtered_df['Revenue'].fillna(0)

        summary = filtered_df.groupby(['Year', 'Month']).agg(
            Case_Volume=('Case Volume', 'sum'),
            Revenue=('Revenue', 'sum')
        ).reset_index()

        summary.rename(columns={'Case_Volume': 'Case Volume', 'Revenue': 'Revenue'}, inplace=True)
        summary['Product Group'] = 'ONLY_ABCD'
        summary['Product Type'] = 'ONLY_ABCD'
        result_df = pd.concat([self.revenue_data, summary], ignore_index=True)
        result_df.drop_duplicates()
        return result_df

    def preprocess_data(self):

        self.revenue_data.drop_duplicates(inplace=True)
        data_with_only_abcd = self.compute_ONLY_ABCD()
        data_with_only_abcd['Year'] = pd.to_numeric(data_with_only_abcd['Year'], errors='coerce')
        data_with_only_abcd = data_with_only_abcd.groupby(['Product Group', 'Product Type', 'Year', 'Month', 'Revenue'], as_index=False).agg({'Case Volume': 'sum'}).reset_index(drop=False)

        if self.include_support_cost:
            support_cost_with_abcd = self.add_support_cost()
            support_cost_with_abcd['Year'] = pd.to_numeric(support_cost_with_abcd['Year'], errors='coerce')
            df = pd.merge(data_with_only_abcd, support_cost_with_abcd[['Product Group', 'Product Type', 'Year', 'Month', 'Support Cost']], on=['Product Group', 'Product Type', 'Year', 'Month'], how='left')
            df['Support Cost'] = pd.to_numeric(df['Support Cost'], errors='coerce')
            df['Support Cost'] = (df['Support Cost'] / 1000000)
            print(f"Support cost successfully added. Current columns {df.columns}")
        if not self.include_support_cost:
            df = data_with_only_abcd

        print(df.columns)
        df[self.column2] = pd.to_numeric(df[self.column2], errors='coerce')
        df[self.column1] = pd.to_numeric(df[self.column1], errors='coerce')

        df['time'] = df['Year'].astype(str) + '-' + df['Month'].astype(str) + '-01'
        df['time'] = pd.to_datetime(df['time'])
        df.set_index('time', inplace=True)
        df.sort_values(by='time', inplace=True)
        df[self.column1] = df[self.column1].astype(float)
        df.loc[:, self.column1] = df[self.column1].round(2)
        df.loc[:, self.column2] = df[self.column2].round(2)
        df.to_csv('df_preprocess.csv')
        return df

    def calculate_statistics(self, df):
        ### Calculate Correlation
        df.fillna(0, inplace=True)
        correlations = df.groupby('Product Group').apply(self.compute_correlation).reset_index()
        return correlations

    def compute_correlation(self, group):
        group.fillna(0, inplace=True)

        stddev_case_volume = group[self.column1].std()
        stddev_revenue = group[self.column2].std()

        if stddev_case_volume == 0 or stddev_revenue == 0:
            #print(f"Zero standard deviation detected for group: {group['Product Group'].iloc[0]}")
            correlation = np.nan
        else:
            correlation = group[self.column1].corr(group[self.column2])

        return pd.Series({
            'Correlation': correlation,
            'Stddev Case Volume': stddev_case_volume,
            'Stddev Revenue': stddev_revenue
        })

    def create_combined_graph(self, df):

        grouped_df = df.groupby(['Product Group', 'time'])[self.column1].sum().reset_index()

        revenue_grouped_df = df.groupby(['Product Group', 'time'])[self.column2].sum().reset_index()
        ### plotting variables

        current_dir = os.path.dirname(os.path.realpath(__file__))
        save_path = os.path.join(current_dir, 'image_bin')
        if self.include_support_cost:
            color='orange'
        else:
            color='blue'

        for product_type in df['Product Group'].unique():

            product_df = grouped_df[grouped_df['Product Group'] == product_type]

            plt.figure(figsize=(10, 15))
            plt.plot(product_df['time'], product_df[self.column1], marker='x', linestyle='dotted', color=color,
                     label=self.column1)
            for i, txt in enumerate(product_df[self.column1]):
                plt.text(product_df['time'].iloc[i], product_df[self.column1].iloc[i], f"{txt}", fontsize=8,
                         ha='right',
                         va='bottom')
            plt.title(f'{self.column1} Time Series for {product_type}')
            plt.xlabel('Year and Quarter')
            plt.ylabel(self.column1)
            plt.xticks(rotation=45)
            plt.legend()
            plt.tight_layout()

            save_path_case_volume = os.path.join(save_path, f'{product_type}_case_volume.png')
            plt.savefig(save_path_case_volume)
            plt.close()

            product_df = revenue_grouped_df[revenue_grouped_df['Product Group'] == product_type]
            product_df.loc[:, self.column2] = product_df[self.column2].round(2)

            plt.figure(figsize=(15, 6))
            plt.plot(product_df['time'], product_df[self.column2], marker='x', linestyle='-', color='green',
                     label=f'{self.column2} ($M)')
            for i, txt in enumerate(product_df[self.column2]):
                plt.text(product_df['time'].iloc[i], product_df[self.column2].iloc[i], f"{txt}", fontsize=8, ha='right',
                         va='bottom')
            plt.title(f'{self.column2} Time Series for {product_type}')
            plt.xlabel('Year and Quarter')
            plt.ylabel(f'{self.column2} ($M)')
            plt.xticks(rotation=45)
            plt.legend()

            plt.tight_layout()
            save_path_revenue = os.path.join(save_path, f'{product_type}_revenue.png')
            plt.savefig(save_path_revenue)
            # plt.show()
            plt.close()

        ############################### COMBINED GRAPHS ##############################

        for product_type in df['Product Group'].unique():

            product_df = grouped_df[grouped_df['Product Group'] == product_type]
            revenue_df = revenue_grouped_df[revenue_grouped_df['Product Group'] == product_type]

            fig, ax1 = plt.subplots(figsize=(15, 6))  ###graph size adjustment

            # Plot Case Volume
            ax1.plot(product_df['time'], product_df[self.column1], marker='x', linestyle='-', color=color,
                     label=self.column1)
            for i, txt in enumerate(product_df[self.column1]):
                ax1.text(product_df['time'].iloc[i], product_df[self.column1].iloc[i], f"{txt}", fontsize=8,
                         ha='right',
                         va='bottom')
            ax1.set_xlabel('Year and Month')
            ax1.set_ylabel(f'{self.column1}', color=color)
            ax1.tick_params(axis='y', colors=color)

            # Add trend line for Case Volume FOR MONTHLY DATA
            z2 = np.polyfit(product_df['time'].map(dt.datetime.toordinal), product_df[self.column1], 1)
            p2 = np.poly1d(z2)
            ax1.plot(product_df['time'], p2(product_df['time'].map(dt.datetime.toordinal)), linestyle='--', color=color,
                     label=f"Trend line - {self.column1}")
            #### red trend line for Case Volume
            last_three_months = product_df['time'] >= product_df['time'].max() - pd.DateOffset(months=3)
            ax1.plot(product_df['time'][last_three_months],
                     p2(product_df['time'][last_three_months].map(dt.datetime.toordinal)), "r--",
                     label="Trend line - Last Quarter")

            # Set x-axis labels with rotation and monthly format
            months = mdates.MonthLocator()  # every month
            months_fmt = mdates.DateFormatter('%Y-%m')
            ax1.xaxis.set_major_locator(months)
            ax1.xaxis.set_major_formatter(months_fmt)
            ax1.xaxis.set_ticklabels(product_df['time'].dt.strftime('%Y-%m'), rotation=45)

            product_df = revenue_grouped_df[revenue_grouped_df['Product Group'] == product_type]

            product_df[self.column2] = product_df[self.column2].round(2)  #### ROUNDING IS HERE

            # Plot Revenue
            ax2 = ax1.twinx()
            ax2.plot(product_df['time'], product_df[self.column2], marker='x', linestyle='-', color='green',
                     label=f'{self.column2} ($M)')

            for i, txt in enumerate(product_df[self.column2]):
                ax2.text(product_df['time'].iloc[i], product_df[self.column2].iloc[i], f"{txt}", fontsize=8, ha='right',
                         va='bottom')
            ax2.set_ylabel(f'{self.column2} ($M)', color='green')
            ax2.tick_params(axis='y', colors='green')

            # Add trend line for Revenue for MONTHLY DATA
            z = np.polyfit(product_df['time'].map(dt.datetime.toordinal), product_df[self.column2], 1)
            p = np.poly1d(z)
            ax2.plot(product_df['time'], p(product_df['time'].map(dt.datetime.toordinal)), "g--",
                     label=f"Trend line - {self.column2}")

            # Red trend line for Revenue (last three months)
            last_three_months = product_df['time'] >= product_df['time'].max() - pd.DateOffset(months=3)
            ax2.plot(product_df['time'][last_three_months],
                     p(product_df['time'][last_three_months].map(dt.datetime.toordinal)), "r--",
                     label="Trend line - Last Quarter")

            # Add title and legend
            fig.suptitle(f'{self.column1} and {self.column2} Time Series for {product_type}')
            ax1.legend()
            ax2.legend()

            # Add title and legend
            fig.suptitle(f'{self.column1} and {self.column2} Time Series for {product_type}')
            ax1.legend(loc='upper left')  # move the first legend to the upper left corner
            ax2.legend(loc='lower right')  # move the second legend to the upper right corner


            plt.tight_layout()
            save_path_combined = os.path.join(save_path, f'{product_type}_case_volume_revenue.png')
            plt.savefig(save_path_combined)
            # plt.show()
            # plt.close()

    def create_wide_summary(self, data):
        data['Year'] = data['Year'].astype(str)
        data['Product Group'] = data['Product Group'].astype(str)
        data = data[data['Product Group'] != 'ONLY_ABCD']
        data[self.column1] = pd.to_numeric(data[self.column1], errors='coerce')
        data[self.column2] = pd.to_numeric(data[self.column2], errors='coerce')

        summary = data.groupby(['Year', 'Product Group']).agg({self.column1: 'sum', self.column2: 'sum'}).reset_index()
        summary[f'{self.column1}/{self.column2} Ratio'] = (
                    summary[self.column1] / summary[self.column2].replace(0, pd.NA))
        wide_summary = summary.pivot(index='Product Group', columns='Year')
        wide_summary.columns = [' '.join(col).strip() for col in wide_summary.columns.values]
        years = summary['Year'].unique()
        metrics = [self.column1, self.column2, f'{self.column1}/{self.column2} Ratio']
        new_column_order = [f'{metric} {year}' for year in years for metric in metrics]
        wide_summary = wide_summary[new_column_order]

        grand_totals = {}
        for year in years:
            total_volume = wide_summary[f'{self.column1} {year}'].sum()
            total_revenue = wide_summary[f'{self.column2} {year}'].sum()
            volume_revenue_ratio = np.nan if total_revenue == 0 else total_volume / total_revenue
            grand_totals[f'{self.column1} {year}'] = total_volume
            grand_totals[f'{self.column2} {year}'] = total_revenue
            grand_totals[f'{self.column1}/{self.column2} Ratio {year}'] = volume_revenue_ratio
        wide_summary.loc['Grand Total'] = pd.Series(grand_totals)

        for year in years:
            wide_summary.loc[:, f'{self.column1} {year}'] = wide_summary[f'{self.column1} {year}'].round(2)
            wide_summary.loc[:, f'{self.column2} {year}'] = wide_summary[f'{self.column2} {year}'].round(2)

        pd.options.display.float_format = '{:,.2f}'.format

        if f'{self.column1} 2022' in wide_summary:
            wide_summary = wide_summary.dropna(subset=[f'{self.column1} 2022'])

        dfi.export(wide_summary, 'wide_summary.png')

    def add_support_cost(self):
        support_cost_full = self.support_cost
        product_categorization = self.product_categorization
        if 'Product Family' in support_cost_full.columns:
            support_cost_full.drop(['Product Family'], axis=1, inplace=True)

        df_merged = pd.merge(support_cost_full,
                             product_categorization[['Product Name', 'Product Group', 'Product Type']],
                             left_on='Product', right_on='Product Name', how='inner')
        print(f"Columns from df_merged while adding support cost {df_merged['Product Group'].unique()}")

        df_merged.fillna(value='NAN', inplace=True)
        df_merged['Product Group'] = df_merged['Product Group'].str.upper()
        df_merged['Product Type'] = df_merged['Product Type'].str.upper()
        df_merged['Product Support Cost'] = df_merged['Product Support Cost'].replace({'\$': '', ',': ''}, regex=True)
        df_merged['Product Support Cost'] = pd.to_numeric(df_merged['Product Support Cost'], errors='coerce')
        df_merged['Month Cases'] = pd.to_numeric(df_merged['Month Cases'], errors='coerce')

        df_merged.dropna(subset=['Product Support Cost', 'cases'], inplace=True)
        df_merged['New Support Cost'] = df_merged['Product Support Cost'] * df_merged['cases']
        df_merged.to_csv('merged.csv')
        long_df = df_merged

        month_mapping = {
            'January': 1,
            'February': 2,
            'March': 3,
            'April': 4,
            'May': 5,
            'June': 6,
            'July': 7,
            'August': 8,
            'September': 9,
            'October': 10,
            'November': 11,
            'December': 12
        }
        subgroups_to_summarize = ["Product_A", "Product_B", "Product_C", "Product_D"]

        long_df['Month'] = long_df['Month of Date'].replace(month_mapping, inplace=False)
        long_df.reset_index(drop=True, inplace=True)

        # Adding only_ABCD
        if self.switch_product_group:
            long_df['Product Group'] = long_df['Product Group'].str.replace('-', ' ')
            filtered_subgroups = long_df[long_df['Product Group'].isin(subgroups_to_summarize)]

        else:
            long_df['Product Type'] = long_df['Product Type'].str.replace('-', ' ')

            filtered_subgroups = long_df[long_df['Product Type'].isin(subgroups_to_summarize)]

        summarized_df = filtered_subgroups.groupby(['Month', 'Year of Date'])['New Support Cost'].sum().reset_index()

        summarized_df['Product Group'] = 'ONLY_ABCD'
        summarized_df['Product Type'] = 'ONLY_ABCD'
        summarized_df = summarized_df[
            ['Product Group', 'Product Type', 'Month', 'Year of Date', 'New Support Cost']]

        summarized_df.to_csv('summarized.csv')
        long_df = long_df[['Product Group', 'Product Type', 'Month', 'Year of Date', 'New Support Cost']]
        long_df = pd.concat([long_df, summarized_df], ignore_index=True)
        long_df.rename(columns={'New Support cost': 'Support Cost'}, inplace=True)
        df_grouped = long_df.groupby(['Product Group', 'Product Type', 'Month', 'Year of Date'])[
            'New Support Cost'].sum().reset_index()
        df_grouped['Year of Date'] = df_grouped['Year of Date'].astype(str)

        df_grouped.rename(columns={'New Support Cost': 'Support Cost'}, inplace=True)
        df_grouped.rename(columns={'Year of Date': 'Year'}, inplace=True)
        df_grouped.rename(columns={'Product Type': 'Product Type'}, inplace=True)
        df_grouped.to_csv('df_grouped_support_cost.csv')
        return df_grouped

    def one_big_graph(self):
        df = self.preprocess_data()
        df = df[df['Product Group'] != 'ONLY_ABCD']
        xe_cogs = pd.read_csv('xe_cogs_spent.csv')
        self.column1 = 'cogs_rounded'
        # Group data by 'Year' and 'Month' and aggregate the sums for 'Revenue' and 'Case Volume'
        df_grouped = df.groupby(['Year', 'Month']).agg({
            self.column2: 'sum'  # Summing Case Volume
        }).reset_index()
        df_grouped['Year'] = df_grouped['Year'].astype(int)
        df_grouped['Month'] = df_grouped['Month'].astype(int)
        xe_cogs['Year'] = xe_cogs['Year'].astype(int)
        xe_cogs['Month'] = xe_cogs['Month'].astype(int)

        df_grouped = pd.merge(df_grouped, xe_cogs, on=['Year', 'Month'])

        df_grouped['Year-Month'] = df_grouped['Year'].astype(str) + '-' + df_grouped['Month'].astype(str).str.zfill(2)

        # Create the plot
        fig, ax1 = plt.subplots(figsize=(12, 6))  # Increased width to give more space

        # Plot Revenue (self.column1) in blue on the primary y-axis
        ax1.set_xlabel('Year-Month')
        ax1.set_ylabel(self.column1, color='tab:blue')
        ax1.plot(df_grouped['Year-Month'], df_grouped[self.column1], color='tab:blue', label=self.column1)

        # Add a trend line for Revenue (self.column1)
        z_revenue = np.polyfit(df_grouped.index, df_grouped[self.column1], 1)
        p_revenue = np.poly1d(z_revenue)
        ax1.plot(df_grouped['Year-Month'], p_revenue(df_grouped.index), linestyle='--', color='tab:blue',
                 label=f"Trend line - {self.column1}")

        # Annotate actual values for Revenue on the plot
        for i, value in enumerate(df_grouped[self.column1]):
            ax1.text(i, value, f"{value:.2f}", color='tab:blue', fontsize=8, ha='right', va='bottom')

        ax1.tick_params(axis='y', labelcolor='tab:blue')

        # Create a secondary y-axis for Case Volume (self.column2) in green
        ax2 = ax1.twinx()
        ax2.set_ylabel(self.column2, color='tab:green')
        ax2.plot(df_grouped['Year-Month'], df_grouped[self.column2], color='tab:green', label=self.column2)

        # Add a trend line for Case Volume (self.column2)
        z_case_volume = np.polyfit(df_grouped.index, df_grouped[self.column2], 1)
        p_case_volume = np.poly1d(z_case_volume)
        ax2.plot(df_grouped['Year-Month'], p_case_volume(df_grouped.index), linestyle='--', color='tab:green',
                 label=f"Trend line - {self.column2}")

        # Annotate actual values for Case Volume on the plot
        for i, value in enumerate(df_grouped[self.column2]):
            ax2.text(i, value, f"{value:.2f}", color='tab:green', fontsize=8, ha='right', va='bottom')

        ax2.tick_params(axis='y', labelcolor='tab:green')

        # Rotate x-axis labels for better readability
        ax1.set_xticks(range(len(df_grouped['Year-Month'])))
        ax1.set_xticklabels(df_grouped['Year-Month'], rotation=45, ha='right')

        # Add title
        plt.title(f'{self.column1} and {self.column2} Time Series', fontsize=14)

        # Show legend with adjusted positions to avoid overlap
        ax1.legend(loc='upper left')  # Move Revenue legend outside the plot
        ax2.legend(loc='lower right')  # Move Case Volume legend outside the plot

        # Save the plot
        plt.savefig('BIG_ONE.png')

        # Show the plot

    def create_presentation(self):
        df = self.preprocess_data()
        print(df.columns)
        df.to_csv('df.csv')
        df['Product Group'] = df['Product Group'].str.replace("/", " ", regex=False)
        self.create_combined_graph(df)
        self.create_wide_summary(df)

        ### PRESENTATION CREATION

        prs = Presentation()

        ###intro slide
        intro_slide_layout = prs.slide_layouts[0]
        intro_slide = prs.slides.add_slide(intro_slide_layout)
        intro_title = intro_slide.shapes.title
        intro_title.text = "Revenue / Support Cases Relationship"
        current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        tx_box = intro_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Presentation generated on: {current_time}"
        correlations = self.calculate_statistics(df)

        # do the thing
        for product_type in df['Product Group'].unique():
            #print(f"Going through product {product_type}")
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            product_df = df[df['Product Group'] == product_type]
            correlation = correlations[correlations['Product Group'] == product_type]['Correlation'].values[0]

            ### Adding correlation into each slide's title
            title = slide.shapes.title
            title.text = f"{product_type} (r = {correlation:.2f})"

            ### Adding plots
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            img_width = Inches(6)
            img_height = Inches(4)
            left = (slide_width - img_width) / 2
            top = (slide_height - img_height) / 2

            image_path_case_volume_revenue = os.path.join(self.image_dir, f'{product_type}_case_volume_revenue.png')
            #slide.shapes.add_picture(image_path_case_volume_revenue, Inches(5.5), Inches(1.5), width=Inches(4.5), height=Inches(3))
            slide.shapes.add_picture(image_path_case_volume_revenue, left, top, width=img_width, height=img_height)

        ### adding wide_summary png table
        table_slide_layout = prs.slide_layouts[5]
        table_slide = prs.slides.add_slide(table_slide_layout)
        title = table_slide.shapes.title
        title.text = f"New metric suggestion - {self.column1} / {self.column2} ratio"
        img_path = 'wide_summary.png'
        table_slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
        prs.save('product_analysis_presentation.pptx')
        print("Presentation created and saved.")


